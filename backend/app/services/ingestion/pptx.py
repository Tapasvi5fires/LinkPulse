import logging
from typing import List
from datetime import datetime
import os
from app.services.ingestion.base import BaseIngestor, IngestedDocument
from pptx import Presentation

logger = logging.getLogger(__name__)

class PptxIngestor(BaseIngestor):
    async def ingest(self, source: str, **kwargs) -> List[IngestedDocument]:
        """
        Ingest a PPTX file from a file path.
        """
        from starlette.concurrency import run_in_threadpool
        return await run_in_threadpool(self._ingest_sync, source)

    def _ingest_sync(self, source: str) -> List[IngestedDocument]:
        documents = []
        try:
            prs = Presentation(source)
            full_text = []

            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    full_text.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))

            text_content = "\n\n".join(full_text)

            if text_content.strip():
                documents.append(IngestedDocument(
                    content=text_content,
                    source_url=source,
                    source_type="pptx",
                    metadata={
                        "title": os.path.basename(source),
                        "slide_count": len(prs.slides),
                        "ingested_at": datetime.utcnow().isoformat(),
                        "source_url": source, # Inject for persistence
                        "source_type": "pptx"
                    }
                ))
        except Exception as e:
            logger.error(f"Error ingesting PPTX {source}: {e}")
            raise e

        return documents
